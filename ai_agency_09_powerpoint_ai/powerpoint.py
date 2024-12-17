from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from openai import OpenAI
from dotenv import load_dotenv  
import os
import json

load_dotenv()

def get_ai_content(topic):
    """
    Get presentation content from OpenAI
    """
    client = OpenAI()
    
    prompt = f"""Create a presentation outline about {topic} with:
    1. A compelling title
    2. 5 main sections
    3. 3-4 bullet points for each section
    Format as JSON with structure:
    {{
        "title": "main title",
        "sections": [
            {{
                "title": "section title",
                "points": ["point 1", "point 2", "point 3"]
            }}
        ]
    }}"""
    
    response = client.chat.completions.create(
        model="gpt-4-turbo-preview",
        messages=[{"role": "user", "content": prompt}],
        temperature=0.7
    )
    
    # Extract JSON content from the response
    content = response.choices[0].message.content
    json_content = content.strip().lstrip('```json').rstrip('```')
    
    # Parse the JSON content
    return json.loads(json_content)

def create_styled_presentation(content, output_file="presentation.pptx"):
    """
    Create a beautifully styled PowerPoint presentation
    """
    prs = Presentation()
    
    # Define color scheme
    COLORS = {
        'primary': RGBColor(33, 60, 114),    # Deep blue
        'secondary': RGBColor(103, 140, 177), # Light blue
        'accent': RGBColor(255, 127, 80),     # Coral
        'background': RGBColor(240, 240, 240) # Light gray
    }
    
    # Title slide
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    
    # Style title slide
    title = title_slide.shapes.title
    subtitle = title_slide.placeholders[1]
    
    title.text = content['title']
    subtitle.text = "Created with AI Assistance"
    
    # Style title text
    title_frame = title.text_frame
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    title_frame.paragraphs[0].font.size = Pt(44)
    title_frame.paragraphs[0].font.color.rgb = COLORS['primary']
    title_frame.paragraphs[0].font.bold = True
    
    # Style subtitle text
    subtitle_frame = subtitle.text_frame
    subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    subtitle_frame.paragraphs[0].font.size = Pt(24)
    subtitle_frame.paragraphs[0].font.color.rgb = COLORS['secondary']
    subtitle_frame.paragraphs[0].font.italic = True
    
    # Content slides
    for section in content['sections']:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        
        # Add and style title
        title = slide.shapes.title
        title.text = section['title']
        title_frame = title.text_frame
        title_frame.paragraphs[0].font.size = Pt(36)
        title_frame.paragraphs[0].font.color.rgb = COLORS['primary']
        title_frame.paragraphs[0].font.bold = True
        
        # Add and style content
        body = slide.shapes.placeholders[1]
        tf = body.text_frame
        tf.word_wrap = True
        
        # Clear default text
        tf.clear()
        
        # Add points with custom styling
        for idx, point in enumerate(section['points']):
            p = tf.add_paragraph()
            p.text = point
            p.font.size = Pt(24)
            p.font.color.rgb = COLORS['secondary']
            p.space_before = Pt(12)
            p.space_after = Pt(12)
            p.level = 0
            
            # Add different color for first point
            if idx == 0:
                p.font.color.rgb = COLORS['accent']
    
    # Save the presentation
    prs.save(output_file)
    return f"Presentation saved as {output_file}"

# Example usage
def main():
    # Set your OpenAI API key as environment variable first
    # os.environ['OPENAI_API_KEY'] = 'your-api-key-here'
    
    # Get AI-generated content
    topic = "The Python Programming Language"
    content = get_ai_content(topic)
    
    # Create the presentation
    result = create_styled_presentation(content)
    print(result)

if __name__ == "__main__":
    main()
