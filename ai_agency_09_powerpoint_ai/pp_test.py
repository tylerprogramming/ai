from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from openai import OpenAI
from pydantic import BaseModel
from typing import List

class Section(BaseModel):
    text: str
    rgb_color: str
    color: str

class SlideContent(BaseModel):
    sections: List[Section]

def get_section_content():
    client = OpenAI()
    
    completion = client.beta.chat.completions.parse(
        model="gpt-4o",
        messages=[
            {
                "role": "system", 
                "content": "You are a presentation design expert. Create content for three sections of AI Agents. Each section should have two lines of text (3-4 words per line and if need lines, add /n and a complementary hex color code."
            },
            {
                "role": "user", 
                "content": "Generate content for three sections, create simple sections."
            }
        ],
        response_format=SlideContent,
    )
    
    content = completion.choices[0].message.parsed

    print(content)
    
    # Convert hex colors to RGB tuples
    for section in content.sections:
        try:
            hex_color = section.color.lstrip('#')
            print(f"hex_color: {hex_color}")
            if len(hex_color) == 6:
                r, g, b = tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))

                section.rgb_color = RGBColor(r, g, b)
            else:
                section.rgb_color = RGBColor(100, 100, 100)  # Black
        except ValueError:
            section.rgb_color = RGBColor(100, 100, 100)  # Black
    
    return content.sections

def create_python_intro_slide(output_file="python_introduction.pptx"):
    # Create presentation and slide
    prs = Presentation()
    
    # Set slide dimensions to 16:9 aspect ratio
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Add blank slide
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    
    # Add dark blue background for title
    title_bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        left=Inches(0),
        top=Inches(0),
        width=prs.slide_width,
        height=Inches(2)
    )
    title_bg.fill.solid()
    title_bg.fill.fore_color.rgb = RGBColor(28, 37, 65)  # Dark blue color
    
    # Add title shape
    title = slide.shapes.add_textbox(
        left=Inches(0),
        top=Inches(0),
        width=prs.slide_width,
        height=Inches(2)
    )
    
    # Style title
    title_tf = title.text_frame
    title_tf.text = "Introduction to Python"
    title_para = title_tf.paragraphs[0]
    title_para.alignment = PP_ALIGN.CENTER
    title_para.font.size = Pt(60)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(255, 255, 255)
    
    # Calculate total width needed for all sections
    num_sections = 3
    section_width = Inches(3.5)  # Width of each section
    total_width = section_width * num_sections
    
    # Calculate starting x position to center all sections
    space_between = Inches(1)  # Space between sections
    total_width_with_spacing = total_width + (space_between * (num_sections - 1))
    start_x = (prs.slide_width - total_width_with_spacing) / 2
    
    def add_section(index, icon_shape, text, icon_color):
        # Calculate position for this section
        left = start_x + (index * (section_width + space_between))
        
        # Center position for icon within section
        icon_width = Inches(1.5)
        icon_left = left + ((section_width - icon_width) / 2)
        
        # Add shape for icon
        icon = slide.shapes.add_shape(
            icon_shape,
            left=icon_left,
            top=Inches(3),
            width=icon_width,
            height=Inches(1.5)
        )
        
        # Style icon
        icon.fill.solid()
        icon.fill.fore_color.rgb = icon_color
        icon.line.color.rgb = icon_color
        
        # Add textbox for the description
        textbox = slide.shapes.add_textbox(
            left=left,
            top=Inches(5),
            width=section_width,
            height=Inches(1)
        )
        
        # Style text
        tf = textbox.text_frame
        tf.clear()  # Clear any existing text
        
        # Split the text into lines and add each as a separate paragraph
        lines = text.split('\n')
        for i, line in enumerate(lines):
            p = tf.add_paragraph()
            p.text = line
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(18)
            p.font.color.rgb = RGBColor(0, 0, 0)
        
        return icon, textbox
    
    # Get AI-generated content for sections
    try:
        sections_content = get_section_content()

        print("sections_content")
        print(sections_content)
        
        # Define sections with their icons and AI-generated content
        sections = [
            {
                'shape': MSO_SHAPE.CUBE,
                'color': sections_content[0].rgb_color,
                'text': sections_content[0].text
            },
            {
                'shape': MSO_SHAPE.ROUNDED_RECTANGLE,
                'color': sections_content[1].rgb_color,
                'text': sections_content[1].text
            },
            {
                'shape': MSO_SHAPE.RIGHT_ARROW,
                'color': sections_content[2].rgb_color,
                'text': sections_content[2].text
            }
        ]

        print("sections")
        print(sections)
    except Exception as e:
        print(f"Error generating content: {e}")
        # Fallback to default content if API call fails
        sections = [
            {
                'shape': MSO_SHAPE.CUBE,
                'color': RGBColor(176, 80, 80),
                'text': "History and philosophy\nbehind Python"
            },
            {
                'shape': MSO_SHAPE.ROUNDED_RECTANGLE,
                'color': RGBColor(134, 176, 80),
                'text': "Python's position in the\nprogramming world today"
            },
            {
                'shape': MSO_SHAPE.RIGHT_ARROW,
                'color': RGBColor(128, 96, 176),
                'text': "Key features that set\nPython apart"
            }
        ]
    
    # Add all sections
    section_shapes = []
    for i, section in enumerate(sections):
        icon, text = add_section(
            index=i,
            icon_shape=section['shape'],
            text=section['text'],
            icon_color=section['color']
        )
        section_shapes.append((icon, text))
    
    # Save presentation
    prs.save(output_file)
    return f"Presentation saved as {output_file}"

# Create the slide
if __name__ == "__main__":
    result = create_python_intro_slide()
    print(result)
